"""Render a resolved pik layout (:mod:`pypik.pik.layout`) to a PowerPoint file.

Coordinates in a :class:`~pypik.pik.layout.LayoutResult` are inches with
y pointing up (pikchr's convention); PowerPoint slides use EMU with y
pointing down from the top-left, so this module flips y and adds a margin
around the diagram's bounding box.

Shape-kind to PowerPoint mapping is necessarily approximate for a few
classes (cylinder, file) since python-pptx's autoshape set doesn't have
an exact equivalent; see the per-kind comments below.
"""

from __future__ import annotations

import os

from PIL import ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

from .pik import ast
from .pik.layout import LayoutResult, Shape, _font_scale, assign_text_slots, resolve_layout

EMU_PER_INCH = 914400

# Vertical offset (in line-steps, positive = up in pik's y-up space) for
# each text slot pik_txt_vertical_layout() can assign -- see assign_text_slots().
_SLOT_STEP = {"above2": 2, "above": 1, "center": 0, "below": -1, "below2": -2}

# The font *name* embedded in the .pptx (what the viewer's PowerPoint/
# Keynote/LibreOffice actually renders with -- Arial is close to always
# available) can differ from the font *file* used here to measure text
# width for "fit" sizing, since that requires an actual font file on this
# machine and Arial itself may not be installed here. PilFontMetrics tries
# a few widely-available metrically-similar substitutes, in order.
#
# A Latin-only substitute (Liberation/DejaVu/Arial) silently *undersizes*
# any CJK text: missing-glyph fallback advances are far narrower than a
# real ideograph, so "fit" shapes come out too small and the text overflows
# them once PowerPoint actually renders it with a CJK-capable font. A Noto
# Sans CJK file, where present, covers Latin *and* CJK correctly, so it's
# tried first; the Latin-only substitutes remain as a fallback chain for
# machines without it (where only non-CJK text will still measure well).
FONT_NAME = "Arial"
_MEASURE_FONT_CANDIDATES = [
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc",
    "/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttf",
    "/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc",
    "/System/Library/Fonts/Supplemental/NotoSansCJKjp-Regular.otf",
    "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "/Library/Fonts/Arial.ttf",
    "/System/Library/Fonts/Supplemental/Arial.ttf",
    "C:\\Windows\\Fonts\\arial.ttf",
]

_BASE_FONT_PT = 9
_LABEL_STEP_IN = 0.10  # how far a label shifts per above/below slot step at _BASE_FONT_PT -- see _SLOT_STEP


def _find_measure_font() -> str | None:
    for path in _MEASURE_FONT_CANDIDATES:
        if os.path.exists(path):
            return path
    return None


class PilFontMetrics:
    """A pypik.pik.layout.FontMetrics backed by real glyph widths (via
    Pillow), so a "fit" object's size tracks its actual text content --
    unlike a fixed per-character-width guess, this gets more (not less)
    accurate as text gets longer or more varied.

    The font used for *measurement* is whatever metrically-close substitute
    for `font_name` this machine actually has installed (see
    _MEASURE_FONT_CANDIDATES); the .pptx itself still requests `font_name`,
    so on a viewer that has it installed, rendering and measurement agree
    exactly. Where they don't exactly agree, this still generalizes
    correctly across different text content, unlike a flat estimate.
    """

    def __init__(self, font_name: str = FONT_NAME, base_size_pt: float = _BASE_FONT_PT, font_path: str | None = None):
        self.font_name = font_name
        self._base_size_pt = base_size_pt
        self._font_path = font_path if font_path is not None else _find_measure_font()
        self._cache: dict[int, ImageFont.FreeTypeFont] = {}

    def _font(self, flags: list[str]) -> tuple[ImageFont.FreeTypeFont | None, float]:
        size_pt = self._base_size_pt * _font_scale(flags)
        size_px = max(1, round(size_pt))
        if size_px not in self._cache:
            self._cache[size_px] = ImageFont.truetype(self._font_path, size_px) if self._font_path else None
        return self._cache[size_px], size_pt

    def text_width(self, text: str, flags: list[str] = ()) -> float:
        font, size_pt = self._font(flags)
        if font is not None:
            px = font.getlength(text) if text else 0.0
        else:
            # No real font file found on this machine at all: fall back to
            # a flat estimate rather than failing outright.
            px = len(text) * size_pt * 0.55
        return px / 72.0 + (size_pt * 0.5) / 72.0

    def line_height(self, flags: list[str] = ()) -> float:
        _font, size_pt = self._font(flags)
        return size_pt / 72.0


def resolve_for_pptx(doc: ast.Document, font_name: str = FONT_NAME, base_size_pt: float = _BASE_FONT_PT) -> LayoutResult:
    """resolve_layout(), using real font metrics so "fit" objects are
    sized to match what write_pptx() will actually draw."""
    return resolve_layout(doc, metrics=PilFontMetrics(font_name, base_size_pt))


def _font_size(flags: list[str], base_size_pt: float = _BASE_FONT_PT) -> Pt:
    return Pt(base_size_pt * _font_scale(flags))


_AUTOSHAPE = {
    "box": MSO_SHAPE.RECTANGLE,
    "circle": MSO_SHAPE.OVAL,
    "ellipse": MSO_SHAPE.OVAL,
    "oval": MSO_SHAPE.OVAL,
    "diamond": MSO_SHAPE.DIAMOND,
    "cylinder": MSO_SHAPE.CAN,  # closest built-in equivalent
    "file": MSO_SHAPE.FLOWCHART_DOCUMENT,  # closest built-in equivalent
    "dot": MSO_SHAPE.OVAL,
    "text": MSO_SHAPE.RECTANGLE,  # rendered with no fill/outline, see below
}


def _rgb(value: float) -> RGBColor:
    v = int(value) & 0xFFFFFF
    return RGBColor((v >> 16) & 0xFF, (v >> 8) & 0xFF, v & 0xFF)


_MIN_SLIDE_SIDE_IN = 1.0  # PowerPoint refuses a slide smaller than 1in on either side


class _Transform:
    """Maps pik inches (y-up) to slide inches (y-down), with a margin.

    PowerPoint requires each slide dimension to be at least 1in; a small
    diagram (or a very small margin) can fall under that on its own, so
    slide_width/slide_height are clamped up to it, and the extra space is
    split evenly as additional margin so the diagram stays centered
    rather than pinned in a corner.
    """

    def __init__(self, bbox: tuple[float, float, float, float], margin: float):
        self.x0, self.y0, self.x1, self.y1 = bbox
        natural_w = (self.x1 - self.x0) + 2 * margin
        natural_h = (self.y1 - self.y0) + 2 * margin
        self.slide_width = max(natural_w, _MIN_SLIDE_SIDE_IN)
        self.slide_height = max(natural_h, _MIN_SLIDE_SIDE_IN)
        self.margin_x = margin + (self.slide_width - natural_w) / 2
        self.margin_y = margin + (self.slide_height - natural_h) / 2

    def rect(self, shape: Shape) -> tuple[float, float, float, float]:
        """Return (left, top, width, height) in inches for shape's bbox."""
        bx0, by0, bx1, by1 = shape.bbox
        left = (bx0 - self.x0) + self.margin_x
        top = (self.y1 - by1) + self.margin_y
        return left, top, bx1 - bx0, by1 - by0

    def point(self, pt: tuple[float, float]) -> tuple[float, float]:
        return (pt[0] - self.x0) + self.margin_x, (self.y1 - pt[1]) + self.margin_y


def _set_arrowheads(line, larrow: bool, rarrow: bool) -> None:
    """python-pptx has no high-level arrowhead API; add the OOXML elements
    directly. <a:headEnd> is the line's start, <a:tailEnd> its end."""
    ln = line._get_or_add_ln()
    for tag, present in (("a:headEnd", larrow), ("a:tailEnd", rarrow)):
        el = ln.find(qn(tag))
        if el is None:
            el = ln.makeelement(qn(tag), {})
            ln.append(el)
        el.set("type", "triangle" if present else "none")


def _apply_line_style(line, shape: Shape) -> None:
    if shape.sw < 0:
        line.fill.background()
        return
    line.color.rgb = _rgb(shape.color)
    line.width = Pt(max(shape.sw, 0.001) * 72)
    if shape.dashed > 0:
        line.dash_style = MSO_LINE_DASH_STYLE.DASH
    elif shape.dotted > 0:
        line.dash_style = MSO_LINE_DASH_STYLE.ROUND_DOT


def _apply_text(pptx_shape, shape: Shape, font_name: str, base_size_pt: float) -> None:
    if not shape.texts:
        return
    tf = pptx_shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # PowerPoint's default autoshape text margins (~0.1in sides, ~0.05in
    # top/bottom) eat into the box width _autosize_text() already sized to
    # the text, which can force unwanted wrapping. pikchr's own charWidth
    # padding is the only margin that's meant to apply, and that's already
    # folded into _autosize_text()'s width/height estimate.
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    for i, (text, flags) in enumerate(shape.texts):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = (
            PP_ALIGN.LEFT if "ljust" in flags else PP_ALIGN.RIGHT if "rjust" in flags else PP_ALIGN.CENTER
        )
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.bold = "bold" in flags
        run.font.italic = "italic" in flags
        run.font.size = _font_size(flags, base_size_pt)
        run.font.color.rgb = _rgb(shape.color)


def _add_block_shape(slide, shape: Shape, tf: _Transform, font_name: str, base_size_pt: float) -> None:
    left, top, w, h = tf.rect(shape)
    w, h = max(w, 0.01), max(h, 0.01)
    autoshape_type = _AUTOSHAPE.get(shape.kind, MSO_SHAPE.RECTANGLE)
    if shape.kind == "box" and shape.rad > 0:
        autoshape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    pptx_shape = slide.shapes.add_shape(autoshape_type, Inches(left), Inches(top), Inches(w), Inches(h))
    if autoshape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
        # `adjustments[0]` is the corner radius as a fraction of min(w, h),
        # not an absolute length.
        pptx_shape.adjustments[0] = min(0.5, shape.rad / min(w, h))

    if shape.kind == "text":
        pptx_shape.fill.background()
        pptx_shape.line.fill.background()
    else:
        if shape.fill < 0:
            pptx_shape.fill.background()
        else:
            pptx_shape.fill.solid()
            pptx_shape.fill.fore_color.rgb = _rgb(shape.fill)
        _apply_line_style(pptx_shape.line, shape)

    _apply_text(pptx_shape, shape, font_name, base_size_pt)


def _add_line_shape(slide, shape: Shape, tf: _Transform, font_name: str, base_size_pt: float) -> None:
    assert shape.path is not None
    points = [tf.point(p) for p in shape.path]

    if len(points) == 2:
        (x1, y1), (x2, y2) = points
        connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        _apply_line_style(connector.line, shape)
        _set_arrowheads(connector.line, shape.larrow, shape.rarrow)
    else:
        x0, y0 = points[0]
        builder = slide.shapes.build_freeform(Inches(x0), Inches(y0))
        builder.add_line_segments([(Inches(x), Inches(y)) for x, y in points[1:]], close=shape.closed)
        freeform = builder.convert_to_shape()
        freeform.fill.background()
        _apply_line_style(freeform.line, shape)
        _set_arrowheads(freeform.line, shape.larrow, shape.rarrow)

    _add_line_text(slide, shape, tf, font_name, base_size_pt)


def _line_label_rects(
    shape: Shape, base_size_pt: float
) -> list[tuple[tuple[float, float, float, float], str, list[str]]]:
    """Compute each text label's (x0, y0, x1, y1) rect in pik space (y-up,
    unmargined), alongside its text/flags -- shared by _add_line_text()
    (which draws these) and _content_bbox() (which needs to know how far
    they extend beyond shape.bbox, since a line's own bbox -- just its
    path -- doesn't account for labels floating above/below it)."""
    if not shape.texts:
        return []
    label_box_h = base_size_pt * 1.15 / 72.0  # a touch taller than line_height(), just for rendering safety
    label_step = base_size_pt / 9.0 * _LABEL_STEP_IN
    bx0, by0, bx1, by1 = shape.bbox
    cx = (bx0 + bx1) / 2
    cy = (by0 + by1) / 2
    out = []
    for (text, flags), slot in zip(shape.texts, assign_text_slots(shape.texts)):
        dy = _SLOT_STEP.get(slot, 0) * label_step
        box_w = max(len(text) * base_size_pt / 72.0 * 0.7, 0.3)
        y = cy + dy
        out.append(((cx - box_w / 2, y - label_box_h / 2, cx + box_w / 2, y + label_box_h / 2), text, flags))
    return out


def _add_line_text(slide, shape: Shape, tf: _Transform, font_name: str, base_size_pt: float) -> None:
    """A connector/freeform shape has no text_frame in python-pptx, so a
    line's text (e.g. an arrow's label) is rendered as small floating
    textboxes instead, placed above/on/below the line per assign_text_slots()."""
    for (x0, y0, x1, y1), text, flags in _line_label_rects(shape, base_size_pt):
        left, top = tf.point((x0, y1))
        textbox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(x1 - x0), Inches(y1 - y0))
        text_frame = textbox.text_frame
        text_frame.word_wrap = False
        text_frame.margin_left = text_frame.margin_right = 0
        text_frame.margin_top = text_frame.margin_bottom = 0
        para = text_frame.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        run = para.add_run()
        run.text = text
        run.font.name = font_name
        run.font.bold = "bold" in flags
        run.font.italic = "italic" in flags
        run.font.size = _font_size(flags, base_size_pt)
        run.font.color.rgb = _rgb(shape.color)


def _content_bbox(result: LayoutResult, base_size_pt: float) -> tuple[float, float, float, float]:
    """result.bbox, expanded to also cover line labels -- a line's own
    bbox is just its path, so a label floating above/below it (see
    _line_label_rects()) can stick out past result.bbox on its own."""
    x0, y0, x1, y1 = result.bbox
    for shape in result.shapes:
        if shape.kind not in ("line", "arrow", "spline", "arc"):
            continue
        for (lx0, ly0, lx1, ly1), _text, _flags in _line_label_rects(shape, base_size_pt):
            x0, y0 = min(x0, lx0), min(y0, ly0)
            x1, y1 = max(x1, lx1), max(y1, ly1)
    return x0, y0, x1, y1


def write_pptx(
    result: LayoutResult,
    path: str,
    margin: float = 0.15,
    font_name: str = FONT_NAME,
    base_size_pt: float = _BASE_FONT_PT,
) -> None:
    """Render `result` (from :func:`pypik.pik.layout.resolve_layout`, or
    `resolve_for_pptx`) to a single-slide PowerPoint file at `path`, sized
    to fit the diagram. `font_name`/`base_size_pt` should match whatever
    was passed to `resolve_for_pptx()` for that result, if it was used, so
    rendering and "fit" sizing agree.

    `margin` only needs to cover the diagram's own edge (e.g. a thick
    stroke's outer half, or PowerPoint's arrowhead overshoot) -- line
    labels are already accounted for by _content_bbox(), not by margin."""
    tf = _Transform(_content_bbox(result, base_size_pt), margin)
    prs = Presentation()
    prs.slide_width = Emu(int(tf.slide_width * EMU_PER_INCH))
    prs.slide_height = Emu(int(tf.slide_height * EMU_PER_INCH))
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

    for shape in result.shapes:
        if shape.kind in ("line", "arrow", "spline", "arc"):
            _add_line_shape(slide, shape, tf, font_name, base_size_pt)
        else:
            _add_block_shape(slide, shape, tf, font_name, base_size_pt)

    prs.save(path)
