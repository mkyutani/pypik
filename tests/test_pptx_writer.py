"""Tests for pypik.pptx_writer: rendering a resolved layout to PowerPoint."""

from __future__ import annotations

import pathlib

import pytest
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_FILL_TYPE
from pptx.enum.shapes import MSO_SHAPE_TYPE

from pypik.pik import parse
from pypik.pptx_writer import resolve_for_pptx, write_pptx

FIXTURES_DIR = pathlib.Path(__file__).parent / "fixtures" / "examples"
EXAMPLE_FILES = sorted(FIXTURES_DIR.glob("*.pik"))


def render(text: str, tmp_path: pathlib.Path, name: str = "out.pptx") -> Presentation:
    # Uses resolve_for_pptx() (real font-metrics-based "fit" sizing), the
    # same path the CLI takes, so these tests exercise it too.
    result = resolve_for_pptx(parse(text))
    out = tmp_path / name
    write_pptx(result, str(out))
    return Presentation(str(out))


@pytest.mark.parametrize("path", EXAMPLE_FILES, ids=lambda p: p.name)
def test_official_example_renders_without_error(path: pathlib.Path, tmp_path: pathlib.Path):
    prs = render(path.read_text(encoding="utf-8"), tmp_path, path.stem + ".pptx")
    assert len(prs.slides) == 1


def test_box_becomes_a_shape_with_text(tmp_path: pathlib.Path):
    prs = render('box "Hello"\n', tmp_path)
    slide = prs.slides[0]
    assert len(slide.shapes) == 1
    shape = slide.shapes[0]
    assert shape.text_frame.text == "Hello"


def test_two_point_arrow_becomes_a_connector(tmp_path: pathlib.Path):
    prs = render("arrow right\n", tmp_path)
    slide = prs.slides[0]
    assert len(slide.shapes) == 1
    assert slide.shapes[0].shape_type == MSO_SHAPE_TYPE.LINE


def test_arrow_labels_render_as_textboxes_above_and_below(tmp_path: pathlib.Path):
    # A connector shape has no text_frame of its own in python-pptx, so an
    # arrow's text must show up as separate floating textboxes -- and two
    # un-flagged texts split above/below the line (pik_txt_vertical_layout()).
    prs = render('arrow right "Top" "Bottom"\n', tmp_path)
    slide = prs.slides[0]
    textboxes = [s for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.TEXT_BOX]
    assert {tb.text_frame.text for tb in textboxes} == {"Top", "Bottom"}
    top_box = next(tb for tb in textboxes if tb.text_frame.text == "Top")
    bottom_box = next(tb for tb in textboxes if tb.text_frame.text == "Bottom")
    # slide y grows downward, so the "above the line" label has the smaller top.
    assert top_box.top < bottom_box.top


def test_rounded_box_gets_rounded_rectangle_with_scaled_corner(tmp_path: pathlib.Path):
    prs = render("box rad 0.1 width 1 height 1\n", tmp_path)
    shape = prs.slides[0].shapes[0]
    assert shape.adjustments[0] == pytest.approx(0.1, abs=0.01)


def test_multipoint_line_becomes_a_freeform(tmp_path: pathlib.Path):
    prs = render("line right 1 then up 1\n", tmp_path)
    slide = prs.slides[0]
    assert len(slide.shapes) == 1
    assert slide.shapes[0].shape_type == MSO_SHAPE_TYPE.FREEFORM


def test_slide_size_matches_diagram_bbox_plus_margin(tmp_path: pathlib.Path):
    result = resolve_for_pptx(parse("box\n"))
    out = tmp_path / "sized.pptx"
    write_pptx(result, str(out), margin=0.5)
    prs = Presentation(str(out))
    x0, y0, x1, y1 = result.bbox
    assert prs.slide_width.inches == pytest.approx((x1 - x0) + 1.0, abs=0.01)
    assert prs.slide_height.inches == pytest.approx((y1 - y0) + 1.0, abs=0.01)


def test_fill_color_is_applied(tmp_path: pathlib.Path):
    # Color names must be capitalized in pikchr source (e.g. "Red") -- that
    # is what makes the tokenizer treat them as a PLACENAME/color reference
    # rather than a lowercase variable name.
    prs = render("box fill Red\n", tmp_path)
    shape = prs.slides[0].shapes[0]
    assert shape.fill.fore_color.rgb == RGBColor(0xFF, 0x00, 0x00)


def test_invis_object_has_no_outline(tmp_path: pathlib.Path):
    prs = render("box invis\n", tmp_path)
    shape = prs.slides[0].shapes[0]
    assert shape.line.fill.type == MSO_FILL_TYPE.BACKGROUND
